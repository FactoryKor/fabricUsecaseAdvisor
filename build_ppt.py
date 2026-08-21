#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Fabric UseCase Advisor 소개 PPT를 python-pptx로 생성한다. 실행: python build_ppt.py"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
BLUE = RGBColor(0x3A, 0x7B, 0xD9)
LIGHT_BLUE = RGBColor(0xEE, 0xF4, 0xFF)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x2E, 0xA8, 0x4A)
ORANGE = RGBColor(0xE8, 0xB4, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, font="Malgun Gothic", anchor=None, line_spacing=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = font
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=DARK, font="Malgun Gothic",
                 line_spacing=1.3):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, lvl = item
        else:
            text, lvl = item, 0
        p.text = ("•  " if lvl == 0 else "-  ") + text
        p.level = lvl
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size if lvl == 0 else size - 2)
            run.font.color.rgb = color
            run.font.name = font
    return box


def add_box(slide, left, top, width, height, text, fill_color=BLUE, text_color=WHITE, size=14,
             bold=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = text_color
        run.font.name = "Malgun Gothic"
    return shape


def add_arrow(slide, left, top, width, height, color=GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    shape.shadow.inherit = False
    return shape


def add_footer(slide, page_no):
    add_text(slide, 0.5, 7.1, 8, 0.3, "Fabric 데이터 활용 시나리오 추천 & 갭 분석", size=10, color=GRAY)
    add_text(slide, 12.3, 7.1, 0.6, 0.3, str(page_no), size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def add_table(slide, left, top, width, height, rows, header_color=BLUE, col_widths=None):
    n_rows, n_cols = len(rows), len(rows[0])
    gtable = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                     Inches(width), Inches(height)).table
    if col_widths:
        for i, w in enumerate(col_widths):
            gtable.columns[i].width = Inches(w)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = gtable.cell(r, c)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(12 if r == 0 else 11)
                    run.font.name = "Malgun Gothic"
                    run.font.bold = (r == 0)
                    run.font.color.rgb = WHITE if r == 0 else DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color if r == 0 else (LIGHT_BLUE if r % 2 == 0 else WHITE)
    return gtable


# ─────────────────────────────────────────────────────────────────────────
# 1. 타이틀
# ─────────────────────────────────────────────────────────────────────────
s = add_slide()
set_background(s, NAVY)
add_text(s, 1, 2.5, 11.3, 1.2, "Fabric 데이터 활용 시나리오 추천 & 갭 분석",
         size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 1, 3.7, 11.3, 0.8, "데이터 소스 → 유스케이스 추천 → 갭 분석 자동화 파이프라인",
         size=20, color=RGBColor(0xCC, 0xDD, 0xF5), align=PP_ALIGN.CENTER)
add_text(s, 1, 6.3, 11.3, 0.5, "Fabric-UseCase-Advisor  |  Prototype",
         size=14, color=RGBColor(0x88, 0xA6, 0xD0), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────
# 2. 문제 정의
# ─────────────────────────────────────────────────────────────────────────
s = add_slide()
add_text(s, 0.6, 0.4, 12, 0.8, "왜 필요한가", size=30, bold=True, color=NAVY)
add_box(s, 0.6, 1.4, 12.1, 0.05, "", fill_color=BLUE, shape_type=MSO_SHAPE.RECTANGLE)
add_bullets(s, 0.8, 1.8, 11.7, 4.5, [
    "기업은 이미 많은 데이터(OneLake, DB, Databricks, SAP 등)를 보유하고 있지만,",
    ("\"이 데이터로 무엇을 할 수 있는지\" 파악하는 과정은 여전히 수작업/컨설팅에 의존한다.", 1),
    "데이터 카탈로그 도구는 \"이 컬럼이 무엇인지\"는 설명해도 \"무슨 업무에 쓰라\"는 추천까지는 하지 못한다.",
    "추천을 하더라도 \"지금 데이터로 부족한 부분이 무엇인지(갭)\"까지 정량적으로 짚어주는 도구는 거의 없다.",
    "결과적으로 신규 유스케이스 발굴에 오랜 시간과 전문 인력이 소요된다.",
], size=18, line_spacing=1.5)
add_footer(s, 2)

# ─────────────────────────────────────────────────────────────────────────
# 3. 핵심 아이디어 — 5단계 파이프라인
# ─────────────────────────────────────────────────────────────────────────
s = add_slide()
add_text(s, 0.6, 0.4, 12, 0.8, "핵심 아이디어: 5단계 자동화 파이프라인", size=28, bold=True, color=NAVY)
add_box(s, 0.6, 1.3, 12.1, 0.05, "", fill_color=BLUE, shape_type=MSO_SHAPE.RECTANGLE)

steps = [
    ("1. 데이터 프로파일링", "스키마·null비율·카디널리티·\n시계열/PII 자동 스캔"),
    ("2. 비즈니스 개념 매핑", "규칙 기반 + LLM으로\n컬럼 → 의미 부여"),
    ("3. 시나리오 추천", "업무 템플릿과 매칭해\n스코어링 (규칙+LLM 하이브리드)"),
    ("4. 갭 분석", "부족한 데이터를\n구체적으로 제안"),
    ("5. 대시보드", "추천 카드 +\n체크리스트 시각화"),
]
box_w, gap, arrow_w = 2.05, 0.15, 0.35
x = 0.6
top = 2.3
for i, (title, desc) in enumerate(steps):
    add_box(s, x, top, box_w, 1.6, title, fill_color=BLUE, text_color=WHITE, size=15)
    add_text(s, x - 0.1, top + 1.7, box_w + 0.2, 1.2, desc, size=11, color=GRAY, align=PP_ALIGN.CENTER)
    x += box_w
    if i < len(steps) - 1:
        add_arrow(s, x, top + 0.55, arrow_w, 0.5, color=ORANGE)
        x += arrow_w
    x += 0.02

add_text(s, 0.8, 5.6, 11.5, 1.2,
         "→ pandas/Spark 프로파일링 + 규칙 기반 매칭을 기본 골격으로 하고, 애매한 구간(컬럼 의미 추정,\n"
         "자연어 설명, 추가 시나리오 브레인스토밍)에 Azure OpenAI(LLM)를 하이브리드로 결합한다.",
         size=15, color=DARK, line_spacing=1.4)
add_footer(s, 3)

# ─────────────────────────────────────────────────────────────────────────
# 4. 차별화 포인트
# ─────────────────────────────────────────────────────────────────────────
s = add_slide()
add_text(s, 0.6, 0.4, 12, 0.8, "차별화 포인트", size=30, bold=True, color=NAVY)
add_box(s, 0.6, 1.4, 12.1, 0.05, "", fill_color=BLUE, shape_type=MSO_SHAPE.RECTANGLE)

diffs = [
    ("① Fabric/OneLake 네이티브 통합", "경쟁 툴 대부분은 범용 카탈로그라 Fabric 특화 기능이 없음 — 이 솔루션은 OneLake/Lakehouse를 1급 데이터 소스로 취급"),
    ("② \"추천\"에서 그치지 않는 정량적 갭 분석", "시나리오별 필요 피처셋 대비 보유 데이터를 diff해 \"무엇이 부족한지\" 구체적으로 제시 — 시중에 거의 없는 기능"),
    ("③ LLM 기반 자연어 설명", "비개발자 현업 담당자도 이해할 수 있는 문장으로 실현 가능성과 보강 방안을 설명"),
    ("④ 이종 데이터 소스 통합", "OneLake뿐 아니라 PostgreSQL/MySQL/MSSQL, SAP HANA(RISE 포함), Databricks까지 하나의 파이프라인에서 프로파일링"),
]
top = 1.9
for title, desc in diffs:
    add_box(s, 0.7, top, 0.5, 0.5, "", fill_color=ORANGE, shape_type=MSO_SHAPE.OVAL)
    add_text(s, 1.4, top, 10.8, 0.5, title, size=17, bold=True, color=NAVY)
    add_text(s, 1.4, top + 0.5, 10.8, 0.7, desc, size=13, color=GRAY, line_spacing=1.3)
    top += 1.25
add_footer(s, 4)

# ─────────────────────────────────────────────────────────────────────────
# 5. 경쟁 제품 비교
# ─────────────────────────────────────────────────────────────────────────
s = add_slide()
add_text(s, 0.6, 0.4, 12, 0.8, "인접 시장 대비 포지셔닝", size=28, bold=True, color=NAVY)
add_box(s, 0.6, 1.3, 12.1, 0.05, "", fill_color=BLUE, shape_type=MSO_SHAPE.RECTANGLE)
rows = [
    ["구분", "대표 제품", "강점", "약점(우리가 채우는 부분)"],
    ["데이터 카탈로그 + AI", "Microsoft Purview, Alation, Collibra", "데이터 설명·거버넌스 우수", "\"무슨 업무에 쓸지\" 추천은 약함"],
    ["Feature Store", "Databricks Feature Store, Tecton", "피처 재사용에 최적화", "신규 유스케이스 발굴 기능 없음"],
    ["Fabric Copilot", "Microsoft Fabric 내장 Copilot", "코드/쿼리 생성 지원", "전략적 유스케이스 추천은 안 함"],
    ["컨설팅 서비스", "Accenture, Deloitte 등", "심층적인 수동 분석", "자동화 도구 아님, 시간·비용 큼"],
    ["스타트업", "Secoda, Metaphor Data 등", "카탈로그+AI 추천 일부 시도", "Fabric 특화 아님"],
]
add_table(s, 0.6, 1.7, 12.1, 4.7, rows, col_widths=[2.0, 3.0, 3.4, 3.7])
add_footer(s, 5)

# ─────────────────────────────────────────────────────────────────────────
# 6. 아키텍처
# ─────────────────────────────────────────────────────────────────────────
s = add_slide()
add_text(s, 0.6, 0.4, 12, 0.8, "아키텍처 개요", size=30, bold=True, color=NAVY)
add_box(s, 0.6, 1.3, 12.1, 0.05, "", fill_color=BLUE, shape_type=MSO_SHAPE.RECTANGLE)

# 데이터 소스 열
sources = ["OneLake /\nLakehouse", "PostgreSQL /\nMySQL / MSSQL", "SAP HANA\n(RISE 포함)", "Databricks"]
sy = 1.9
for i, src in enumerate(sources):
    add_box(s, 0.7, sy + i * 1.05, 2.3, 0.85, src, fill_color=NAVY, text_color=WHITE, size=12)
add_arrow(s, 3.15, 3.6, 0.55, 0.6, color=ORANGE)

add_box(s, 3.85, 2.3, 3.0, 2.6, "프로파일링 +\n비즈니스 개념 매핑\n(규칙 기반 + LLM)", fill_color=BLUE, text_color=WHITE, size=14)
add_arrow(s, 7.0, 3.6, 0.55, 0.6, color=ORANGE)

add_box(s, 7.7, 2.3, 3.0, 2.6, "시나리오 스코어링 +\n갭 분석\n(규칙 + LLM 하이브리드)", fill_color=BLUE, text_color=WHITE, size=14)
add_arrow(s, 10.85, 3.6, 0.55, 0.6, color=ORANGE)

add_box(s, 11.55, 2.3, 1.2, 2.6, "대시보드\n/\nPower BI", fill_color=GREEN, text_color=WHITE, size=13)

add_text(s, 0.7, 5.3, 12, 1.5,
         "Azure OpenAI(LLM)는 컬럼 의미 추정, 시나리오 자연어 설명, 추가 시나리오 브레인스토밍 단계에서\n"
         "선택적으로 호출되며(모델은 사용자가 배포 이름으로 직접 선택), 미설정 시 규칙 기반만으로도 전체 파이프라인이 동작한다.",
         size=14, color=GRAY, line_spacing=1.4)
add_footer(s, 6)

# ─────────────────────────────────────────────────────────────────────────
# 7. 지원 데이터 소스
# ─────────────────────────────────────────────────────────────────────────
s = add_slide()
add_text(s, 0.6, 0.4, 12, 0.8, "지원 데이터 소스", size=30, bold=True, color=NAVY)
add_box(s, 0.6, 1.3, 12.1, 0.05, "", fill_color=BLUE, shape_type=MSO_SHAPE.RECTANGLE)
rows = [
    ["소스", "연결 방식", "상태"],
    ["OneLake / Lakehouse (Delta)", "Fabric 내장 Spark 세션으로 실제 연결", "실제 연결 지원 (Fabric 환경)"],
    ["PostgreSQL", "psycopg2", "실제 연결 지원"],
    ["MySQL", "pymysql", "실제 연결 지원"],
    ["SQL Server (MSSQL)", "pyodbc + ODBC Driver 17/18", "실제 연결 지원"],
    ["SAP HANA (RISE with SAP 포함)", "hdbcli (hana_diagnose.py와 동일 방식)", "실제 연결 지원"],
    ["Databricks", "databricks-sql-connector", "커넥터 인터페이스 준비 (스텁 → 실연결 전환 가능)"],
]
add_table(s, 0.6, 1.7, 12.1, 4.5, rows, col_widths=[4.2, 4.6, 3.3])
add_text(s, 0.7, 6.4, 11.8, 0.6,
         "※ 연결 정보가 없거나 실패하면 자동으로 합성 데이터로 대체되어 파이프라인이 항상 끝까지 실행됩니다.",
         size=12, color=GRAY)
add_footer(s, 7)

# ─────────────────────────────────────────────────────────────────────────
# 8. LLM 활용 방식
# ─────────────────────────────────────────────────────────────────────────
s = add_slide()
add_text(s, 0.6, 0.4, 12, 0.8, "LLM(생성형 AI) 활용 방식", size=28, bold=True, color=NAVY)
add_box(s, 0.6, 1.3, 12.1, 0.05, "", fill_color=BLUE, shape_type=MSO_SHAPE.RECTANGLE)
add_bullets(s, 0.8, 1.8, 11.7, 4.5, [
    "모델 선택의 자유: 배포 이름(변수) 하나만 바꾸면 어떤 Azure OpenAI 모델이든 사용 가능",
    ("예: gpt-4o-mini, gpt-5.6-sol 등 보유한 배포를 그대로 지정", 1),
    "① 컬럼 → 비즈니스 개념 매핑 보완: 규칙 기반으로 못 찾은 컬럼만 LLM에게 의미 추정 요청",
    "② 시나리오 실현 가능성 자연어 설명: 커버리지 점수 뒤에 \"왜 그런지, 무엇을 보강하면 되는지\"를 문장으로 생성",
    "③ 추가 시나리오 브레인스토밍: 고정 템플릿 밖에서도 보유 데이터만으로 가능한 새 아이디어 제안",
    "인증: 회사 정책으로 API 키 발급이 막힌 환경도 고려해 Entra ID(Windows 계정) 인증 기본 지원",
    "LLM 미설정 시에도 규칙 기반 엔진만으로 전체 파이프라인이 정상 동작(안전한 폴백)",
], size=16, line_spacing=1.45)
add_footer(s, 8)

# ─────────────────────────────────────────────────────────────────────────
# 9. 데모 결과 예시
# ─────────────────────────────────────────────────────────────────────────
s = add_slide()
add_text(s, 0.6, 0.4, 12, 0.8, "데모 결과 예시 — 시나리오 준비도", size=28, bold=True, color=NAVY)
add_box(s, 0.6, 1.3, 12.1, 0.05, "", fill_color=BLUE, shape_type=MSO_SHAPE.RECTANGLE)

demo = [
    ("수요 예측", 100, GREEN),
    ("배송 지연 예측", 100, GREEN),
    ("재고 최적화", 100, GREEN),
    ("고객 세그멘테이션", 100, GREEN),
    ("고객 이탈 예측", 100, GREEN),
    ("이상 거래 탐지", 66.7, ORANGE),
]
top = 1.9
max_w = 8.5
for name, pct, color in demo:
    add_text(s, 0.7, top, 2.6, 0.5, name, size=14, color=DARK)
    bar_w = max_w * (pct / 100.0)
    add_box(s, 3.4, top + 0.05, max(bar_w, 0.05), 0.4, "", fill_color=color,
             shape_type=MSO_SHAPE.RECTANGLE)
    add_text(s, 3.4 + max_w + 0.15, top, 1.2, 0.5, f"{pct}%", size=14, bold=True, color=color)
    top += 0.72

add_text(s, 0.7, 6.4, 11.8, 0.7,
         "실제 보유 데이터에서 \"이상 거래 탐지\"만 66.7% — transaction_id 개념이 부족해 갭으로 식별됨.",
         size=13, color=GRAY)
add_footer(s, 9)

# ─────────────────────────────────────────────────────────────────────────
# 10. 갭 분석 체크리스트 예시
# ─────────────────────────────────────────────────────────────────────────
s = add_slide()
add_text(s, 0.6, 0.4, 12, 0.8, "갭 분석 체크리스트 예시", size=28, bold=True, color=NAVY)
add_box(s, 0.6, 1.3, 12.1, 0.05, "", fill_color=BLUE, shape_type=MSO_SHAPE.RECTANGLE)
checklist = [
    "[수요 예측] 필요한 핵심 데이터가 모두 확보되어 있습니다.",
    "[고객 이탈 예측] 필요한 핵심 데이터가 모두 확보되어 있습니다.",
    "[배송 지연 예측] 필요한 핵심 데이터가 모두 확보되어 있습니다.",
    "[이상 거래 탐지] 'transaction_id' 데이터가 추가로 필요합니다.",
    "[재고 최적화] 필요한 핵심 데이터가 모두 확보되어 있습니다.",
    "[고객 세그멘테이션] 필요한 핵심 데이터가 모두 확보되어 있습니다.",
]
top = 2.0
for item in checklist:
    is_gap = "추가로 필요" in item
    mark_color = ORANGE if is_gap else GREEN
    mark = "!" if is_gap else "✓"
    add_box(s, 0.7, top, 0.45, 0.45, mark, fill_color=mark_color, text_color=WHITE, size=16,
             shape_type=MSO_SHAPE.OVAL)
    add_text(s, 1.35, top + 0.02, 11.0, 0.5, item, size=15, color=DARK)
    top += 0.7
add_footer(s, 10)

# ─────────────────────────────────────────────────────────────────────────
# 11. 로드맵
# ─────────────────────────────────────────────────────────────────────────
s = add_slide()
add_text(s, 0.6, 0.4, 12, 0.8, "다음 단계 (로드맵)", size=30, bold=True, color=NAVY)
add_box(s, 0.6, 1.3, 12.1, 0.05, "", fill_color=BLUE, shape_type=MSO_SHAPE.RECTANGLE)
add_bullets(s, 0.8, 1.8, 11.7, 4.8, [
    "Fabric Notebook으로 실제 OneLake/Lakehouse 데이터 연결 및 검증",
    "Databricks 실제 연결 전환 (databricks-sql-connector)",
    "자격 증명을 환경변수 → Azure Key Vault / Fabric 데이터 소스 연결 기능으로 이전",
    "결과를 Fabric Lakehouse 테이블로 저장 → Power BI 대시보드로 상시 노출",
    "업무 시나리오 템플릿 라이브러리 확대(업종별 특화 템플릿 추가)",
    "실제 고객 데이터 기반 파일럿 진행 및 정확도/유용성 피드백 수집",
], size=17, line_spacing=1.5)
add_footer(s, 11)

# ─────────────────────────────────────────────────────────────────────────
# 12. 마무리
# ─────────────────────────────────────────────────────────────────────────
s = add_slide()
set_background(s, NAVY)
add_text(s, 1, 2.8, 11.3, 1.0, "Q & A", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 1, 4.0, 11.3, 0.6, "감사합니다", size=20, color=RGBColor(0xCC, 0xDD, 0xF5), align=PP_ALIGN.CENTER)

out_path = "Fabric_UseCase_Advisor_소개.pptx"
prs.save(out_path)
print(f"저장 완료: {out_path}")
